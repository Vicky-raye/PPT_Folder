from pathlib import Path
from typing import List, Dict, Any

def extract_pdf_pages(file_path: Path) -> List[Dict[str, Any]]:
    import fitz  # PyMuPDF
    pages = []
    doc = fitz.open(file_path)
    try:
        for i in range(len(doc)):
            page = doc[i]
            text = page.get_text()
            pages.append({"page_index": i + 1, "page_range": str(i + 1), "text": text.strip() or "(无文本)"})
    finally:
        doc.close()
    return pages


def extract_ppt_pages(file_path: Path) -> List[Dict[str, Any]]:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
        text = "\n".join(parts).strip() if parts else "(无文本)"
        pages.append({"page_index": i + 1, "page_range": str(i + 1), "text": text})
    return pages


def extract_pages(file_path: Path, file_type: str) -> List[Dict[str, Any]]:
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(str(path))
    if file_type.lower() in ("pdf", ".pdf"):
        return extract_pdf_pages(path)
    if file_type.lower() in ("ppt", "pptx", ".ppt", ".pptx"):
        return extract_ppt_pages(path)
    raise ValueError(f"Unsupported file type: {file_type}")
